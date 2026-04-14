"""
Teams Bot 訊息處理器
處理所有來自 Teams 的用戶互動，包括文字訊息和卡片互動
"""

import logging
from typing import Dict, Any, List, Optional
from botbuilder.core import TurnContext
from botbuilder.schema import Activity, ActivityTypes, SuggestedActions

from domain.models.user import UserProfile
from domain.services.todo_service import TodoService
from domain.services.conversation_service import ConversationService
from domain.services.meeting_service import MeetingService
from domain.services.intent_service import IntentService
from application.handlers.bot_command_handler import BotCommandHandler
from application.dtos.bot_dtos import BotInteractionDTO
from config.settings import AppConfig
from presentation.cards.card_builders import (
    TodoCardBuilder,
    HelpCardBuilder,
    MeetingCardBuilder,
    ModelSelectionCardBuilder,
    UploadCardBuilder,
)
from shared.utils.helpers import (
    get_user_email,
    determine_language,
    get_suggested_replies,
)
from shared.exceptions import OpenAIServiceError
from botbuilder.schema import ActionTypes, CardAction, SuggestedActions


class TeamsMessageHandler:
    """Teams 訊息處理器"""

    def __init__(
        self,
        config: AppConfig,
        todo_service: TodoService,
        conversation_service: ConversationService,
        meeting_service: MeetingService,
        intent_service: IntentService,
        command_handler: BotCommandHandler,
    ):
        self.config = config
        self.todo_service = todo_service
        self.conversation_service = conversation_service
        self.meeting_service = meeting_service
        self.intent_service = intent_service
        self.command_handler = command_handler

        # Card builders
        self.todo_card_builder = TodoCardBuilder()
        self.help_card_builder = HelpCardBuilder()
        self.meeting_card_builder = MeetingCardBuilder()
        self.model_card_builder = ModelSelectionCardBuilder()
        self.upload_card_builder = UploadCardBuilder()
        self.logger = logging.getLogger(__name__)

        # 暫存使用者附件（等待確認是否附加到 IT 工單）
        self._pending_attachments: dict = {}
        # 暫存 @t 發送訊息的目標（email -> {"target_email", "sender_name", "timestamp"}）
        self._pending_send_message: dict = {}

    async def handle_message(self, turn_context: TurnContext) -> None:
        """處理 Teams 訊息"""
        try:
            # 提取用戶信息
            user_info = await self._extract_user_info(turn_context)
            snippet = (user_info.message_text or "").strip()
            if len(snippet) > 120:
                snippet = f"{snippet[:117]}..."
            self.logger.info(
                "Teams message received user_mail=%s conversation_id=%s text=\"%s\"",
                user_info.user_mail,
                user_info.conversation_id,
                snippet or "<empty>",
            )

            # 更新用戶會話參考
            await self._update_user_conversation_ref(turn_context, user_info)

            # 處理卡片互動
            if turn_context.activity.value:
                self.logger.info(
                    "Handling adaptive card interaction user_mail=%s conversation_id=%s",
                    user_info.user_mail,
                    user_info.conversation_id,
                )
                await self._handle_card_interaction(turn_context, user_info)
                return

            # Debug: 記錄所有附件的原始資訊，協助排查 Teams 各種附件格式
            raw_atts = turn_context.activity.attachments or []
            if raw_atts:
                for i, a in enumerate(raw_atts):
                    self.logger.info(
                        "RAW attachment[%d] content_type=%s name=%s contentUrl=%s hasContent=%s",
                        i,
                        getattr(a, "content_type", None),
                        getattr(a, "name", None),
                        bool(getattr(a, "content_url", None) or getattr(a, "contentUrl", None)),
                        bool(getattr(a, "content", None)),
                    )

            # 若含有附件：判斷是否有最近 IT 工單，決定附加或 AI 解析
            # 過濾掉 Teams 非檔案附件（mention text/html、adaptive card 等）
            real_attachments = [
                a for a in (turn_context.activity.attachments or [])
                if not (getattr(a, "content_type", "") or "").lower().startswith((
                    "text/html",
                    "application/vnd.microsoft.card",
                ))
            ]
            if real_attachments:
                self.logger.info(
                    "Attempting attachment handling user_mail=%s conversation_id=%s attachment_count=%d",
                    user_info.user_mail,
                    user_info.conversation_id,
                    len(turn_context.activity.attachments or []),
                )

                # 檢查是否有待轉發的 @t 訊息（5 分鐘內）
                pending_msg = self._pending_send_message.get(user_info.user_mail)
                if pending_msg:
                    from datetime import datetime
                    elapsed = (datetime.now() - pending_msg["timestamp"]).total_seconds()
                    if elapsed <= 300:
                        files = self._parse_attachment_files(turn_context)
                        if files:
                            downloaded = await self._download_parsed_files(files)
                            if downloaded:
                                await self._forward_attachment_to_user(
                                    turn_context, user_info, downloaded,
                                )
                                return
                    else:
                        self._pending_send_message.pop(user_info.user_mail, None)

                # 檢查是否有最近 10 分鐘內的 IT 工單
                from core.container import get_container
                from features.it_support.service import ITSupportService
                svc: ITSupportService = get_container().get(ITSupportService)
                recent_gid = svc.get_recent_task_gid(user_info.user_mail)

                if recent_gid:
                    # 有最近的 IT 單 → 暫存附件，發確認卡片讓使用者選擇
                    self._pending_attachments[user_info.user_mail] = {
                        "activity": turn_context.activity,
                        "user_info": user_info,
                    }
                    language = determine_language(user_info.user_mail)
                    confirm_text = {
                        "zh": "您最近有提交 IT 工單，請問要將此檔案附加到該工單嗎？",
                        "en": "You recently submitted an IT ticket. Attach this file to it?",
                        "ja": "最近ITチケットを提出しました。このファイルを添付しますか？",
                    }.get(language, "您最近有提交 IT 工單，請問要將此檔案附加到該工單嗎？")
                    card = {
                        "$schema": "http://adaptivecards.io/schemas/adaptive-card.json",
                        "type": "AdaptiveCard",
                        "version": "1.3",
                        "body": [{"type": "TextBlock", "text": confirm_text, "wrap": True}],
                        "actions": [
                            {"type": "Action.Submit", "title": "✅ 附加到 IT 工單", "data": {"action": "confirmAttachIT"}},
                            {"type": "Action.Submit", "title": "🤖 AI 解析內容", "data": {"action": "skipAttachIT"}},
                        ],
                    }
                    from botbuilder.schema import Attachment as BotAttachment
                    card_attachment = BotAttachment(
                        content_type="application/vnd.microsoft.card.adaptive",
                        content=card,
                    )
                    await turn_context.send_activity(
                        Activity(type=ActivityTypes.message, attachments=[card_attachment])
                    )
                    return

                # 無最近 IT 單 → 下載附件並用 AI 解析
                files = self._parse_attachment_files(turn_context)
                if files:
                    downloaded = await self._download_parsed_files(files)
                    if downloaded:
                        self.logger.info(
                            "Attachment analysis start user_mail=%s count=%d",
                            user_info.user_mail,
                            len(downloaded),
                        )
                        await self._handle_attachment_analysis(
                            turn_context, user_info, downloaded,
                        )
                        return

            # 處理文字訊息
            await self._handle_text_message(turn_context, user_info)

        except Exception:
            user_info = locals().get("user_info")
            self.logger.exception(
                "處理訊息時發生錯誤 user_mail=%s conversation_id=%s",
                getattr(user_info, "user_mail", "unknown"),
                getattr(user_info, "conversation_id", "unknown"),
            )
            await self._send_error_response(turn_context)

    async def _extract_user_info(self, turn_context: TurnContext) -> BotInteractionDTO:
        """提取用戶信息"""
        user_id = turn_context.activity.from_property.id
        user_name = turn_context.activity.from_property.name
        user_mail = await get_user_email(turn_context) or f"{user_id}@unknown.com"
        conversation_id = turn_context.activity.conversation.id

        # Debug 模式處理
        if self.config.debug_mode and self.config.debug_account:
            user_mail = self.config.debug_account

        masked_email = user_mail
        if "@" in user_mail:
            local, domain = user_mail.split("@", 1)
            masked_email = f"{local[:3]}***@{domain}" if len(local) > 3 else f"{local}***@{domain}"
        self.logger.debug(
            "User extracted name=%s user_id=%s mail=%s conversation_id=%s",
            user_name,
            user_id,
            masked_email,
            conversation_id,
        )

        return BotInteractionDTO(
            user_id=user_id,
            user_name=user_name,
            user_mail=user_mail,
            conversation_id=conversation_id,
            message_text=turn_context.activity.text or "",
            card_action=(
                turn_context.activity.value.get("action")
                if turn_context.activity.value
                else None
            ),
        )

    async def _update_user_conversation_ref(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """更新用戶會話參考"""
        conversation_ref = TurnContext.get_conversation_reference(turn_context.activity)
        # 這裡應該透過服務來更新，而不是直接操作全域變數
        # 暫時保持相容性
        from app import user_conversation_refs, user_display_names

        user_conversation_refs[user_info.user_mail] = conversation_ref
        if user_info.user_mail:
            user_display_names[user_info.user_mail] = (
                user_info.user_name or user_display_names.get(user_info.user_mail)
            )

    async def _handle_card_interaction(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """處理卡片互動"""
        card_action = turn_context.activity.value.get("action")

        if card_action == "selectFunction":
            await self._handle_function_selection(turn_context, user_info)
        elif card_action == "addTodo":
            await self._handle_add_todo_card(turn_context, user_info)
        elif card_action == "completeTodos":
            await self._handle_complete_todos(turn_context, user_info)
        elif card_action == "bookRoom":
            await self._handle_book_room(turn_context, user_info)
        elif card_action == "cancelBooking":
            await self._handle_cancel_booking(turn_context, user_info)
        elif card_action == "selectModel":
            await self._handle_model_selection(turn_context, user_info)
        elif card_action == "submitIT":
            await self._handle_submit_it_issue(turn_context, user_info)
        elif card_action == "submitITT":
            await self._handle_submit_itt_issue(turn_context, user_info)
        elif card_action == "uploadOption":
            await self._handle_upload_option(turn_context, user_info)
        elif card_action == "submitBroadcast":
            await self._handle_submit_broadcast(turn_context, user_info)
        elif card_action == "submitSendMessage":
            await self._handle_submit_send_message(turn_context, user_info)
        elif card_action == "submitReplyToIT":
            await self._handle_submit_reply_to_it(turn_context, user_info)
        elif card_action == "confirmAttachIT":
            await self._handle_confirm_attach_it(turn_context, user_info)
        elif card_action == "skipAttachIT":
            await self._handle_skip_attach_it(turn_context, user_info)
        elif card_action == "submitKB":
            await self._handle_submit_kb_query(turn_context, user_info)

    async def _handle_submit_it_issue(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """提交 IT 提單並建立 Asana 任務（背景處理，避免 Teams 逾時）"""
        try:
            form = {
                "summary": turn_context.activity.value.get("summary"),
                "description": turn_context.activity.value.get("description"),
                "category": turn_context.activity.value.get("category"),
                "priority": turn_context.activity.value.get("priority"),
            }

            if not (form.get("description") or "").strip():
                await turn_context.send_activity(
                    Activity(type=ActivityTypes.message, text="❌ 需求/問題說明不得為空")
                )
                return

            # 立即回應使用者，避免 Teams 15 秒逾時
            await turn_context.send_activity(
                Activity(type=ActivityTypes.message, text="⏳ 正在處理您的需求，請稍候...")
            )

            # 取得 conversation reference 供背景推播使用
            conversation_ref = TurnContext.get_conversation_reference(turn_context.activity)

            # 背景執行提單流程
            import asyncio
            asyncio.create_task(
                self._submit_it_issue_background(
                    form, user_info, conversation_ref
                )
            )
        except Exception as e:
            self.logger.exception("啟動 IT 提單背景任務失敗")
            await turn_context.send_activity(
                Activity(type=ActivityTypes.message, text=f"❌ 提交 IT 提單時發生錯誤：{str(e)}")
            )

    async def _submit_it_issue_background(
        self, form: dict, user_info: BotInteractionDTO,
        conversation_ref,
    ) -> None:
        """背景執行 IT 提單流程，完成後透過 proactive message 通知使用者"""
        try:
            from core.container import get_container
            from features.it_support.service import ITSupportService
            import os

            svc: ITSupportService = get_container().get(ITSupportService)
            result = await svc.submit_issue(form, user_info.user_name or "", user_info.user_mail)

            from infrastructure.bot.bot_adapter import CustomBotAdapter
            bot_adapter: CustomBotAdapter = get_container().get(CustomBotAdapter)
            bot_app_id = os.getenv("BOT_APP_ID") or os.getenv("MICROSOFT_APP_ID") or ""

            if result.get("success"):
                msg_text = result.get("message", "✅ 已建立 IT Issue")

                async def send_result(turn_context):
                    await turn_context.send_activity(
                        Activity(type=ActivityTypes.message, text=msg_text)
                    )
                    try:
                        language = determine_language(user_info.user_mail)
                        upload_card = self.upload_card_builder.build_file_upload_options_card(language)
                        await turn_context.send_activity(upload_card)
                    except Exception:
                        pass
                    # 提示使用者可查詢工單狀態
                    tip = "💡 小提示：輸入 **@itls** 可隨時查看您申請的 IT 工單處理進度。"
                    await turn_context.send_activity(
                        Activity(type=ActivityTypes.message, text=tip)
                    )

                await bot_adapter.adapter.continue_conversation(
                    conversation_ref, send_result, bot_app_id
                )
            else:
                error_text = f"❌ {result.get('error', '提交失敗')}"

                async def send_error(turn_context):
                    await turn_context.send_activity(
                        Activity(type=ActivityTypes.message, text=error_text)
                    )

                await bot_adapter.adapter.continue_conversation(
                    conversation_ref, send_error, bot_app_id
                )
        except Exception as e:
            self.logger.exception("IT 提單背景處理失敗: %s", e)
            try:
                from core.container import get_container
                from infrastructure.bot.bot_adapter import CustomBotAdapter
                import os
                bot_adapter: CustomBotAdapter = get_container().get(CustomBotAdapter)
                bot_app_id = os.getenv("BOT_APP_ID") or os.getenv("MICROSOFT_APP_ID") or ""
                err_msg = f"❌ 提交 IT 提單時發生錯誤：{str(e)}"

                async def send_exc(turn_context):
                    await turn_context.send_activity(
                        Activity(type=ActivityTypes.message, text=err_msg)
                    )

                await bot_adapter.adapter.continue_conversation(
                    conversation_ref, send_exc, bot_app_id
                )
            except Exception:
                self.logger.exception("IT 提單背景處理：推播錯誤訊息也失敗")

    async def _handle_submit_itt_issue(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """提交 IT 代提單並建立 Asana 任務（背景處理，避免 Teams 逾時）"""
        try:
            requester_email = (turn_context.activity.value.get("requesterEmail") or "").strip()
            if not requester_email:
                await turn_context.send_activity(
                    Activity(type=ActivityTypes.message, text="❌ 請填寫提出人 Email")
                )
                return

            form = {
                "summary": turn_context.activity.value.get("summary"),
                "description": turn_context.activity.value.get("description"),
                "category": turn_context.activity.value.get("category"),
                "priority": turn_context.activity.value.get("priority"),
            }

            if not (form.get("description") or "").strip():
                await turn_context.send_activity(
                    Activity(type=ActivityTypes.message, text="❌ 需求/問題說明不得為空")
                )
                return

            # 立即回應使用者，避免 Teams 15 秒逾時
            await turn_context.send_activity(
                Activity(type=ActivityTypes.message, text="⏳ 正在處理您的需求，請稍候...")
            )

            conversation_ref = TurnContext.get_conversation_reference(turn_context.activity)

            import asyncio
            asyncio.create_task(
                self._submit_itt_issue_background(
                    form, user_info, conversation_ref, requester_email
                )
            )
        except Exception as e:
            self.logger.exception("啟動 IT 代提單背景任務失敗")
            await turn_context.send_activity(
                Activity(type=ActivityTypes.message, text=f"❌ 提交 IT 代提單時發生錯誤：{str(e)}")
            )

    async def _submit_itt_issue_background(
        self, form: dict, user_info: BotInteractionDTO,
        conversation_ref, requester_email: str,
    ) -> None:
        """背景執行 IT 代提單流程，完成後透過 proactive message 通知使用者"""
        try:
            from core.container import get_container
            from features.it_support.service import ITSupportService
            import os

            svc: ITSupportService = get_container().get(ITSupportService)
            result = await svc.submit_issue(
                form,
                user_info.user_name or "",
                user_info.user_mail,
                requester_email=requester_email,
            )

            from infrastructure.bot.bot_adapter import CustomBotAdapter
            bot_adapter: CustomBotAdapter = get_container().get(CustomBotAdapter)
            bot_app_id = os.getenv("BOT_APP_ID") or os.getenv("MICROSOFT_APP_ID") or ""

            if result.get("success"):
                msg_text = result.get("message", "✅ 已建立 IT Issue（代提單）")

                async def send_result(turn_context):
                    await turn_context.send_activity(
                        Activity(type=ActivityTypes.message, text=msg_text)
                    )
                    try:
                        language = determine_language(user_info.user_mail)
                        upload_card = self.upload_card_builder.build_file_upload_options_card(language)
                        await turn_context.send_activity(upload_card)
                    except Exception:
                        pass
                    tip = "💡 小提示：輸入 **@itls** 可隨時查看您申請的 IT 工單處理進度。"
                    await turn_context.send_activity(
                        Activity(type=ActivityTypes.message, text=tip)
                    )

                await bot_adapter.adapter.continue_conversation(
                    conversation_ref, send_result, bot_app_id
                )
            else:
                error_text = f"❌ {result.get('error', '提交失敗')}"

                async def send_error(turn_context):
                    await turn_context.send_activity(
                        Activity(type=ActivityTypes.message, text=error_text)
                    )

                await bot_adapter.adapter.continue_conversation(
                    conversation_ref, send_error, bot_app_id
                )
        except Exception as e:
            self.logger.exception("IT 代提單背景處理失敗: %s", e)
            try:
                from core.container import get_container
                from infrastructure.bot.bot_adapter import CustomBotAdapter
                import os
                bot_adapter: CustomBotAdapter = get_container().get(CustomBotAdapter)
                bot_app_id = os.getenv("BOT_APP_ID") or os.getenv("MICROSOFT_APP_ID") or ""
                err_msg = f"❌ 提交 IT 代提單時發生錯誤：{str(e)}"

                async def send_exc(turn_context):
                    await turn_context.send_activity(
                        Activity(type=ActivityTypes.message, text=err_msg)
                    )

                await bot_adapter.adapter.continue_conversation(
                    conversation_ref, send_exc, bot_app_id
                )
            except Exception:
                self.logger.exception("IT 代提單背景處理：推播錯誤訊息也失敗")

    async def _handle_function_selection(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """處理功能選擇"""
        selected_function = turn_context.activity.value.get("selectedFunction")
        if not selected_function:
            return

        # 映射功能到處理器
        function_handlers = {
            "@addTodo": self._show_add_todo_card,
            "@ls": self._show_todo_list,
            "@it": self._show_it_issue_card,
            "@itt": self._show_itt_issue_card,
            "@upload": self._show_upload_card,
            "@book-room": self._show_room_booking_options,
            "@check-booking": self._show_my_bookings,
            "@cancel-booking": self._show_cancel_booking_options,
            "@info": self._show_user_info,
            "@you": self._show_bot_intro,
            "@model": self._show_model_selection,
            "@send": self._show_broadcast_card,
        }

        handler = function_handlers.get(selected_function)
        if handler:
            await handler(turn_context, user_info)

    async def _show_it_issue_card(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """顯示 IT 提單卡片"""
        from core.container import get_container
        from features.it_support.service import ITSupportService

        language = determine_language(user_info.user_mail)
        svc: ITSupportService = get_container().get(ITSupportService)
        card = svc.build_issue_card(language, user_info.user_name or "", user_info.user_mail)
        await turn_context.send_activity(card)

    async def _show_itt_issue_card(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """顯示 IT 代提單卡片（含提出人 Email 欄位）"""
        from core.container import get_container
        from features.it_support.service import ITSupportService

        language = determine_language(user_info.user_mail)
        svc: ITSupportService = get_container().get(ITSupportService)
        card = svc.build_itt_issue_card(language, user_info.user_name or "", user_info.user_mail)
        await turn_context.send_activity(card)

    async def _show_broadcast_card(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """顯示廣播推播卡片"""
        from features.it_support.cards import build_broadcast_card
        language = determine_language(user_info.user_mail)
        card = build_broadcast_card(language)
        await turn_context.send_activity(card)

    async def _handle_submit_broadcast(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """處理廣播推播表單提交"""
        try:
            target_emails_raw = (turn_context.activity.value.get("targetEmails") or "").strip()
            message_text = (turn_context.activity.value.get("broadcastMessage") or "").strip()

            if not message_text:
                await turn_context.send_activity(
                    Activity(type=ActivityTypes.message, text="❌ 廣播失敗：推播訊息不能為空")
                )
                return

            if not target_emails_raw:
                await turn_context.send_activity(
                    Activity(type=ActivityTypes.message, text="❌ 廣播失敗：收件人不能為空 (若要全發送請輸入 all)")
                )
                return

            # 從依賴取得 bot_adapter 以及 user_repo
            from core.container import get_container
            from domain.repositories.user_repository import UserRepository
            from presentation.web.api_routes import create_api_routes
            
            container = get_container()
            try:
                user_repo: UserRepository = container.get(UserRepository)
            except Exception as e:
                await turn_context.send_activity(Activity(type=ActivityTypes.message, text=f"❌ 無法獲取用戶庫: {str(e)}"))
                return
            
            # 這裡透過 container.get("bot_adapter") 取出原本注冊好的 Adapter 才能發起 continue_conversation
            # app.py 有註冊 providers.Object(adapter)
            try:
                bot_adapter = container.get("bot_adapter")
            except Exception:
                bot_adapter = None
                
            if not bot_adapter:
                await turn_context.send_activity(Activity(type=ActivityTypes.message, text="❌ 廣播失敗：無法取得 Bot Adapter。"))
                return
                
            bot_app_id = self.config.bot.app_id

            # Parse targets
            target_all = (target_emails_raw.lower() == "all")
            target_list = []
            if not target_all:
                target_list = [t.strip().lower() for t in target_emails_raw.split(";") if t.strip()]

            success_count = 0
            fail_count = 0
            skipped_count = 0
            
            # Send message to targets
            for email, session in user_repo._sessions.items():
                email_lower = email.lower()
                
                # Check if this user is in the target list
                if not target_all and email_lower not in target_list:
                    continue
                    
                ref = session.conversation_reference
                if not ref:
                    skipped_count += 1
                    continue
                    
                try:
                    async def send_proactive_message(turn_context):
                        from botbuilder.schema import Activity
                        activity = Activity(
                            type="message",
                            text=message_text
                        )
                        await turn_context.send_activity(activity)

                    await bot_adapter.adapter.continue_conversation(
                        ref,
                        send_proactive_message,
                        bot_app_id
                    )
                    success_count += 1
                except Exception as e:
                    self.logger.error(f"Failed to send proactive message to {email}: {str(e)}")
                    fail_count += 1

            # Report back
            result_msg = f"✅ 推播完成！\n成功發送：{success_count} 筆\n發送失敗：{fail_count} 筆\n無有效連線略過：{skipped_count} 筆"
            await turn_context.send_activity(
                Activity(type=ActivityTypes.message, text=result_msg)
            )

        except Exception as e:
            self.logger.error(f"Broadcast submission failed: {str(e)}")
            await turn_context.send_activity(
                Activity(type=ActivityTypes.message, text=f"❌ 處理廣播時發生例外：{str(e)}")
            )

    async def _handle_submit_send_message(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """處理 @t 發送訊息給指定使用者"""
        try:
            target_email = (turn_context.activity.value.get("targetUserEmail") or "").strip()
            message_text = (turn_context.activity.value.get("sendMessageText") or "").strip()

            if not target_email:
                await turn_context.send_activity(
                    Activity(type=ActivityTypes.message, text="❌ 請選擇收件人")
                )
                return
            if not message_text:
                await turn_context.send_activity(
                    Activity(type=ActivityTypes.message, text="❌ 訊息內容不能為空")
                )
                return

            from app import user_conversation_refs, user_display_names
            from core.container import get_container
            from infrastructure.bot.bot_adapter import CustomBotAdapter

            ref = user_conversation_refs.get(target_email)
            if not ref:
                await turn_context.send_activity(
                    Activity(type=ActivityTypes.message, text=f"❌ 找不到 {target_email} 的連線資訊，該使用者可能尚未與 Bot 互動。")
                )
                return

            adapter: CustomBotAdapter = get_container().get(CustomBotAdapter)
            bot_app_id = self.config.bot.app_id
            sender_email = user_info.user_mail

            # 發送含回覆按鈕的 Adaptive Card
            from botbuilder.schema import Attachment as BotAttachment
            card_content = {
                "$schema": "http://adaptivecards.io/schemas/adaptive-card.json",
                "type": "AdaptiveCard",
                "version": "1.4",
                "body": [
                    {
                        "type": "TextBlock",
                        "text": "💬 來自 TR GPT 的訊息",
                        "weight": "Bolder",
                        "size": "Medium",
                    },
                    {
                        "type": "TextBlock",
                        "text": message_text,
                        "wrap": True,
                    },
                    {
                        "type": "Input.Text",
                        "id": "replyMessageText",
                        "label": "回覆訊息",
                        "isMultiline": True,
                        "maxLength": 2000,
                        "placeholder": "輸入回覆內容...",
                        "height": "stretch",
                    },
                ],
                "actions": [
                    {
                        "type": "Action.Submit",
                        "title": "↩️ 回覆",
                        "data": {
                            "action": "submitReplyToIT",
                            "replyToEmail": sender_email,
                        },
                    },
                ],
            }

            async def send_card(ctx):
                attachment = BotAttachment(
                    content_type="application/vnd.microsoft.card.adaptive",
                    content=card_content,
                )
                await ctx.send_activity(Activity(type="message", attachments=[attachment]))

            await adapter.adapter.continue_conversation(ref, send_card, bot_app_id)

            # 暫存目標，供後續圖片轉發使用（5 分鐘內有效）
            from datetime import datetime
            self._pending_send_message[user_info.user_mail] = {
                "target_email": target_email,
                "timestamp": datetime.now(),
            }

            target_name = user_display_names.get(target_email, target_email)
            await turn_context.send_activity(
                Activity(
                    type=ActivityTypes.message,
                    text=f"✅ 訊息已成功發送給 {target_name}！\n📎 5 分鐘內可直接貼上或拖曳圖片，我會一併轉發。",
                )
            )

        except Exception as e:
            self.logger.error(f"Send message failed: {str(e)}")
            await turn_context.send_activity(
                Activity(type=ActivityTypes.message, text=f"❌ 發送訊息時發生例外：{str(e)}")
            )

    async def _forward_attachment_to_user(
        self, turn_context: TurnContext, user_info: BotInteractionDTO,
        downloaded: list,
    ) -> None:
        """將附件轉發給 @t 指定的目標使用者"""
        pending = self._pending_send_message.pop(user_info.user_mail, None)
        if not pending:
            return

        target_email = pending["target_email"]

        from app import user_conversation_refs, user_display_names
        from core.container import get_container
        from infrastructure.bot.bot_adapter import CustomBotAdapter
        import base64

        ref = user_conversation_refs.get(target_email)
        if not ref:
            await turn_context.send_activity(
                Activity(type=ActivityTypes.message, text=f"❌ 找不到 {target_email} 的連線資訊。")
            )
            return

        adapter: CustomBotAdapter = get_container().get(CustomBotAdapter)
        bot_app_id = self.config.bot.app_id

        # 將圖片 bytes 轉為 data URI，透過 inline attachment 發送
        from botbuilder.schema import Attachment as BotAttachment
        attachments = []
        for att in downloaded:
            img_bytes = att.get("bytes")
            mime = att.get("mime_type", "application/octet-stream")
            name = att.get("name", "file")
            if img_bytes and mime.startswith("image/"):
                b64 = base64.b64encode(img_bytes).decode("utf-8")
                data_uri = f"data:{mime};base64,{b64}"
                attachments.append(BotAttachment(
                    content_type=mime,
                    content_url=data_uri,
                    name=name,
                ))
            elif img_bytes:
                # 非圖片檔案也用 inline 方式
                b64 = base64.b64encode(img_bytes).decode("utf-8")
                data_uri = f"data:{mime};base64,{b64}"
                attachments.append(BotAttachment(
                    content_type=mime,
                    content_url=data_uri,
                    name=name,
                ))

        if not attachments:
            await turn_context.send_activity(
                Activity(type=ActivityTypes.message, text="⚠️ 無法解析附件內容。")
            )
            return

        async def send_with_attachments(ctx):
            activity = Activity(
                type="message",
                text=f"📎 來自 **TR GPT** 的圖片：",
                attachments=attachments,
            )
            await ctx.send_activity(activity)

        await adapter.adapter.continue_conversation(ref, send_with_attachments, bot_app_id)

        target_name = user_display_names.get(target_email, target_email)
        await turn_context.send_activity(
            Activity(type=ActivityTypes.message, text=f"✅ 圖片已成功轉發給 {target_name}！")
        )

    async def _handle_submit_reply_to_it(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """處理使用者回覆 IT 人員的訊息"""
        try:
            reply_text = (turn_context.activity.value.get("replyMessageText") or "").strip()
            reply_to_email = (turn_context.activity.value.get("replyToEmail") or "").strip()

            if not reply_text:
                await turn_context.send_activity(
                    Activity(type=ActivityTypes.message, text="❌ 回覆內容不能為空")
                )
                return
            if not reply_to_email:
                await turn_context.send_activity(
                    Activity(type=ActivityTypes.message, text="❌ 無法辨識回覆對象")
                )
                return

            from app import user_conversation_refs, user_display_names
            from core.container import get_container
            from infrastructure.bot.bot_adapter import CustomBotAdapter

            ref = user_conversation_refs.get(reply_to_email)
            if not ref:
                await turn_context.send_activity(
                    Activity(type=ActivityTypes.message, text="❌ 該 IT 人員目前不在線上，無法轉發回覆。")
                )
                return

            adapter: CustomBotAdapter = get_container().get(CustomBotAdapter)
            bot_app_id = self.config.bot.app_id
            user_name = user_info.user_name or user_info.user_mail

            async def send_reply(ctx):
                activity = Activity(
                    type="message",
                    text=f"↩️ 來自 **{user_name}** 的回覆：\n\n{reply_text}",
                )
                await ctx.send_activity(activity)

            await adapter.adapter.continue_conversation(ref, send_reply, bot_app_id)

            await turn_context.send_activity(
                Activity(type=ActivityTypes.message, text="✅ 回覆已送出！")
            )

        except Exception as e:
            self.logger.error(f"Reply to IT failed: {str(e)}")
            await turn_context.send_activity(
                Activity(type=ActivityTypes.message, text=f"❌ 回覆時發生例外：{str(e)}")
            )

    async def _show_upload_card(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """顯示檔案上傳引導 HeroCard"""
        language = determine_language(user_info.user_mail)
        # 顯示含 1/2/3 選項的 HeroCard（im_back）
        card = self.upload_card_builder.build_file_upload_options_card(language)
        await turn_context.send_activity(card)

    async def _handle_upload_option(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        opt = str(turn_context.activity.value.get("opt"))
        language = determine_language(user_info.user_mail)
        tips = {
            "zh": {
                "1": "請直接在此對話視窗貼上圖片（或拖曳圖片）後送出，我會自動附加到最近建立的 IT 單。",
                "2": "目前不建議以網址上傳，請改用貼上圖片或 Teams 附件按鈕。",
                "3": "請使用訊息列的附件（迴紋針）按鈕選擇檔案，送出後我會自動附加到最近建立的 IT 單。",
            },
            "en": {
                "1": "Paste or drag the image here and send it; I'll attach it to your latest IT ticket.",
                "2": "URL uploads are not recommended; please paste image or use the attachment button.",
                "3": "Use the attachment (paperclip) button to upload; I'll attach it to your latest IT ticket.",
            },
            "ja": {
                "1": "このチャットに画像を貼り付け／ドラッグして送信してください。最新のITチケットに自動添付します。",
                "2": "URL 経由のアップロードは推奨しません。画像を貼り付けるか、添付ボタンをご利用ください。",
                "3": "メッセージ欄のクリップアイコンからファイルを添付してください。最新のITチケットに自動添付します。",
            },
        }
        text = tips.get(language, tips["zh"]).get(opt)
        if text:
            await turn_context.send_activity(Activity(type=ActivityTypes.message, text=text))

    async def _handle_confirm_attach_it(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """使用者確認將附件附加到 IT 工單"""
        pending = self._pending_attachments.pop(user_info.user_mail, None)
        if not pending:
            await turn_context.send_activity(
                Activity(type=ActivityTypes.message, text="⚠️ 找不到待處理的附件，請重新上傳。")
            )
            return

        # 還原原始 activity 的 attachments，走 IT 工單附加流程
        original_activity = pending["activity"]
        # 建立一個臨時 context 來處理（使用原始 activity 的 attachments）
        turn_context.activity.attachments = original_activity.attachments
        handled = await self._try_attach_images(turn_context, pending["user_info"])
        if not handled:
            await turn_context.send_activity(
                Activity(type=ActivityTypes.message, text="⚠️ 附加失敗，可能工單已超過 10 分鐘。")
            )

    async def _handle_skip_attach_it(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """使用者選擇不附加到 IT 工單，改走 AI 解析"""
        pending = self._pending_attachments.pop(user_info.user_mail, None)
        if not pending:
            await turn_context.send_activity(
                Activity(type=ActivityTypes.message, text="⚠️ 找不到待處理的附件，請重新上傳。")
            )
            return

        # 立即回應，避免 Teams 卡片互動逾時
        await turn_context.send_activity(
            Activity(type=ActivityTypes.message, text="⏳ 正在解析檔案內容，請稍候...")
        )

        # 還原原始 activity 的 attachments，解析檔案清單
        original_activity = pending["activity"]
        turn_context.activity.attachments = original_activity.attachments
        files = self._parse_attachment_files(turn_context)

        # 取得 conversation reference 供背景推播
        conversation_ref = TurnContext.get_conversation_reference(turn_context.activity)

        # 背景執行下載 + AI 解析，避免逾時
        import asyncio
        asyncio.create_task(
            self._skip_attach_it_background(files, pending["user_info"], conversation_ref)
        )

    async def _skip_attach_it_background(
        self, files: list, user_info: BotInteractionDTO, conversation_ref,
    ) -> None:
        """背景執行附件下載與 AI 解析，完成後透過 proactive message 回傳結果"""
        try:
            downloaded = await self._download_parsed_files(files) if files else []
            if not downloaded:
                await self._send_proactive_message(
                    conversation_ref, "⚠️ 無法解析附件，請重新上傳。"
                )
                return

            # 用 proactive message context 執行附件分析
            import os
            from core.container import get_container
            from infrastructure.bot.bot_adapter import CustomBotAdapter

            bot_adapter: CustomBotAdapter = get_container().get(CustomBotAdapter)
            bot_app_id = os.getenv("BOT_APP_ID") or os.getenv("MICROSOFT_APP_ID") or ""

            async def do_analysis(turn_context):
                await self._handle_attachment_analysis(turn_context, user_info, downloaded)

            await bot_adapter.adapter.continue_conversation(
                conversation_ref, do_analysis, bot_app_id
            )
        except Exception as e:
            self.logger.exception("AI 解析附件背景處理失敗: %s", e)
            try:
                await self._send_proactive_message(
                    conversation_ref, "❌ 檔案解析失敗，請稍後再試。"
                )
            except Exception:
                pass

    async def _handle_submit_kb_query(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """處理知識庫查詢卡片提交"""
        kb_slug = (turn_context.activity.value.get("kbSlug") or "").strip()
        question = (turn_context.activity.value.get("kbQuestion") or "").strip()

        if not question:
            await turn_context.send_activity(
                Activity(type=ActivityTypes.message, text="❌ 請輸入查詢問題。")
            )
            return

        if not kb_slug:
            await turn_context.send_activity(
                Activity(type=ActivityTypes.message, text="❌ 請選擇知識庫。")
            )
            return

        # 立即回應避免 Teams 卡片互動逾時
        await turn_context.send_activity(
            Activity(type=ActivityTypes.message, text=f"⏳ 正在查詢知識庫，請稍候...")
        )

        # 背景執行 KB 查詢
        conversation_ref = TurnContext.get_conversation_reference(turn_context.activity)
        import asyncio
        asyncio.create_task(
            self._submit_kb_query_background(kb_slug, question, user_info, conversation_ref)
        )

    async def _submit_kb_query_background(
        self, kb_slug: str, question: str, user_info: BotInteractionDTO, conversation_ref,
    ) -> None:
        """背景執行知識庫查詢，完成後透過 proactive message 回傳結果"""
        try:
            from features.it_support.kb_client import KBVectorClient
            from features.it_support.cards import build_kb_result_card

            kb_client = KBVectorClient()
            result = await kb_client.ask(question, role="kb", kb_name=kb_slug, timeout=30)

            answer = (result.get("answer") or "").strip()
            sources = result.get("sources") or []

            if not answer:
                await self._send_proactive_message(
                    conversation_ref, "⚠️ 知識庫未找到相關資料，請嘗試換個方式描述您的問題。"
                )
                return

            # 用 proactive message 發送結果卡片
            import os
            from core.container import get_container
            from infrastructure.bot.bot_adapter import CustomBotAdapter

            bot_adapter: CustomBotAdapter = get_container().get(CustomBotAdapter)
            bot_app_id = os.getenv("BOT_APP_ID") or os.getenv("MICROSOFT_APP_ID") or ""

            language = determine_language(user_info.user_mail)
            card = build_kb_result_card(language, kb_slug, question, answer, sources)

            async def send_card(turn_context):
                await turn_context.send_activity(card)

            await bot_adapter.adapter.continue_conversation(
                conversation_ref, send_card, bot_app_id
            )
        except Exception as e:
            self.logger.exception("KB 知識庫查詢背景處理失敗: %s", e)
            try:
                await self._send_proactive_message(
                    conversation_ref, "❌ 知識庫查詢失敗，請稍後再試。"
                )
            except Exception:
                pass

    async def _handle_text_message(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """處理文字訊息"""
        user_message = user_info.message_text.strip()

        # 上傳選項快捷回覆（HeroCard im_back）
        if user_message in ["1", "2", "3"]:
            language = determine_language(user_info.user_mail)
            tips = {
                "zh": {
                    "1": "請直接在此對話視窗貼上圖片（或拖曳圖片）後送出，我會自動附加到最近建立的 IT 單。",
                    "2": "目前不建議以網址上傳，請改用貼上圖片或 Teams 附件按鈕。",
                    "3": "請使用訊息列的附件（迴紋針）按鈕選擇檔案，送出後我會自動附加到最近建立的 IT 單。",
                },
                "en": {
                    "1": "Paste or drag the image here and send it; I'll attach it to your latest IT ticket.",
                    "2": "URL uploads are not recommended; please paste image or use the attachment button.",
                    "3": "Use the attachment (paperclip) button to upload; I'll attach it to your latest IT ticket.",
                },
                "ja": {
                    "1": "このチャットに画像を貼り付け／ドラッグして送信してください。最新のITチケットに自動添付します。",
                    "2": "URL 経由のアップロードは推奨しません。画像を貼り付けるか、添付ボタンをご利用ください。",
                    "3": "メッセージ欄のクリップアイコンからファイルを添付してください。最新のITチケットに自動添付します。",
                },
            }
            t = tips.get(language, tips["zh"]).get(user_message)
            if t:
                await turn_context.send_activity(Activity(type=ActivityTypes.message, text=t))
                return

        # 支援 /help 或 help 顯示功能選單（對齊 app_bak 行為）
        if user_message.lower() in ["/help", "help", "@help"]:
            language = determine_language(user_info.user_mail)
            include_model = not self.config.openai.use_azure
            welcome_msg = {
                "zh": "🛠️ 功能選單",
                "en": "🛠️ Function Menu",
                "ja": "🛠️ 機能メニュー",
            }.get(language, "🛠️ 功能選單")
            help_card = self.help_card_builder.build_help_card(
                language, welcome_msg, include_model_option=include_model
            )
            await turn_context.send_activity(help_card)
            return

        # 處理歡迎訊息
        if user_message.lower() in ["hi", "hello", "你好", "嗨"]:
            await self._send_welcome_message(turn_context, user_info)
            return

        # 處理命令
        if user_message.startswith("@"):
            await self.command_handler.handle_command(turn_context, user_info)
            return

        # 處理一般對話（使用 AI 意圖分析或直接 OpenAI 調用）
        if self.config.enable_ai_intent_analysis:
            await self._handle_intent_based_response(turn_context, user_info)
        else:
            await self._handle_direct_openai_response(turn_context, user_info)

    async def _handle_intent_based_response(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """基於意圖分析的回應處理"""
        try:
            intent_result = await self.intent_service.analyze_intent(
                user_info.message_text
            )

            if intent_result.action == "add" and intent_result.category == "todo":
                content = intent_result.content or user_info.message_text
                await self._handle_smart_todo_add(turn_context, user_info, content)
            elif intent_result.action == "query" and intent_result.category == "todo":
                await self._show_todo_list(turn_context, user_info)
            elif intent_result.action == "book" and intent_result.category == "meeting":
                await self._show_room_booking_options(turn_context, user_info)
            elif intent_result.action == "query" and intent_result.category == "meeting":
                await self._show_my_bookings(turn_context, user_info)
            elif intent_result.action == "cancel" and intent_result.category == "meeting":
                await self._show_cancel_booking_options(turn_context, user_info)
            elif intent_result.category == "info":
                if intent_result.action in ("user_info", "status"):
                    await self._show_user_info(turn_context, user_info)
                elif intent_result.action == "bot_info" or intent_result.action == "help":
                    # bot_info/幫助都導向說明/介紹
                    await self._show_bot_intro(turn_context, user_info)
                else:
                    await self._show_user_info(turn_context, user_info)
            elif intent_result.category == "model" and intent_result.action == "select":
                await self._show_model_selection(turn_context, user_info)
            else:
                # 進入主要AI對話 預設回應 由Openai回覆
                # 發送 loading 訊息
                language = determine_language(user_info.user_mail)
                loading_messages = {
                    "zh-TW": "🤔 思考更長時間以取得更佳回答...",
                    "ja": "🤔 考え中です。少々お待ちください...",
                }
                loading_text = loading_messages.get(language, loading_messages["zh-TW"])

                # 發送 typing 活動
                await turn_context.send_activity(Activity(type="typing"))
                await turn_context.send_activity(
                    Activity(type=ActivityTypes.message, text=loading_text)
                )

                await self._handle_direct_openai_response(turn_context, user_info)

        except Exception:
            self.logger.exception(
                "意圖分析失敗 user_mail=%s conversation_id=%s",
                user_info.user_mail,
                user_info.conversation_id,
            )
            await self._handle_direct_openai_response(turn_context, user_info)

    async def _resolve_kb_for_user(self, user_mail: str, question: str) -> str:
        """根據使用者部門 + 個人權限查詢對應知識庫，並行查詢後合併回答。

        兩層映射：
        1. KB_DEPARTMENT_MAP — 部門共用 KB（value 為陣列）
        2. KB_USER_MAP — 個人專屬 KB（機密/主管級）
        合併去重後並行查詢所有 KB，回傳合併的參考文字。
        """
        import os
        import json
        import asyncio

        try:
            from core.container import get_container
            from infrastructure.external.graph_api_client import GraphAPIClient
            from features.it_support.kb_client import KBVectorClient

            container = get_container()
            graph_client: GraphAPIClient = container.get(GraphAPIClient)

            # 取得使用者部門
            user_department = ""
            try:
                user_profile = await graph_client.get_user_info(user_mail)
                if user_profile:
                    user_department = (user_profile.get("department") or "").strip()
            except Exception as e:
                self.logger.warning("⚠️ KB: 無法取得使用者部門: %s", e)

            # --- 第一層：部門共用 KB ---
            dept_map_raw = os.getenv("KB_DEPARTMENT_MAP", "{}")
            try:
                dept_map: dict = json.loads(dept_map_raw)
            except json.JSONDecodeError:
                dept_map = {}

            kb_list: list[str] = []
            if user_department:
                matched = dept_map.get(user_department)
                if not matched:
                    # 模糊匹配
                    for key, slugs in dept_map.items():
                        if key in user_department or user_department in key:
                            matched = slugs
                            break
                if matched:
                    # 相容舊格式（字串）與新格式（陣列）
                    if isinstance(matched, str):
                        kb_list.append(matched)
                    elif isinstance(matched, list):
                        kb_list.extend(matched)

            # --- 第二層：個人專屬 KB ---
            user_map_raw = os.getenv("KB_USER_MAP", "{}")
            try:
                user_map: dict = json.loads(user_map_raw)
            except json.JSONDecodeError:
                user_map = {}

            user_kbs = user_map.get(user_mail) or user_map.get(user_mail.lower()) or []
            if isinstance(user_kbs, str):
                user_kbs = [user_kbs]
            kb_list.extend(user_kbs)

            # 去重並保持順序
            seen = set()
            unique_kbs = []
            for kb in kb_list:
                if kb and kb not in seen:
                    seen.add(kb)
                    unique_kbs.append(kb)

            if not unique_kbs:
                self.logger.debug("KB: user=%s 部門='%s' 無對應知識庫，跳過查詢", user_mail, user_department)
                return ""

            self.logger.info("📚 KB 查詢: user=%s, 部門=%s, KBs=%s", user_mail, user_department, unique_kbs)

            # --- 並行查詢所有 KB（return_exceptions 確保單一 KB 失敗不影響其他） ---
            kb_client = KBVectorClient()
            tasks = [
                kb_client.ask_safe(question, role="user", kb_name=kb)
                for kb in unique_kbs
            ]
            results = await asyncio.gather(*tasks, return_exceptions=True)

            # 合併有效結果，跳過失敗或不存在的 KB
            combined_parts: list[str] = []
            for kb_name, result in zip(unique_kbs, results):
                if isinstance(result, Exception):
                    self.logger.warning("⚠️ KB '%s' 查詢異常（已略過）: %s", kb_name, result)
                    continue
                answer = (result.get("answer") or "").strip()
                if answer and answer != "No relevant knowledge base articles found.":
                    combined_parts.append(f"[{kb_name}] {answer}")
                    self.logger.info("📚 KB 命中: kb=%s, answer_len=%d", kb_name, len(answer))

            if combined_parts:
                return "\n\n".join(combined_parts)

        except Exception as e:
            self.logger.warning("⚠️ KB 查詢失敗（不影響主流程）: %s", e)

        return ""

    async def _handle_direct_openai_response(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """直接 OpenAI 回應處理"""
        # 依模式決定使用的模型（OpenAI 模式支援 per-user 偏好）
        model_arg = None
        if not self.config.openai.use_azure:
            try:
                from app import user_model_preferences
                model_arg = user_model_preferences.get(
                    user_info.user_mail, self.config.openai.model
                )
            except Exception:
                model_arg = self.config.openai.model

        prompt_preview = (user_info.message_text or "").strip()
        if len(prompt_preview) > 120:
            prompt_preview = f"{prompt_preview[:117]}..."

        try:
            self.logger.info(
                "OpenAI response start user_mail=%s conversation_id=%s model=%s prompt=\"%s\"",
                user_info.user_mail,
                user_info.conversation_id,
                model_arg or self.config.openai.model,
                prompt_preview or "<empty>",
            )

            # 查詢使用者部門對應的知識庫
            kb_context = await self._resolve_kb_for_user(user_info.user_mail, user_info.message_text)

            call_kwargs: dict = {}
            if model_arg:
                call_kwargs["model"] = model_arg
            if kb_context:
                call_kwargs["kb_context"] = kb_context

            response = await self.conversation_service.get_ai_response(
                user_info.conversation_id,
                user_info.user_mail,
                user_info.message_text,
                **call_kwargs,
            )

            suggested_actions = get_suggested_replies(
                user_info.message_text, user_info.user_mail
            )

            response_preview = (response or "").strip().replace("\n", " ")
            if len(response_preview) > 120:
                response_preview = f"{response_preview[:117]}..."
            self.logger.info(
                "OpenAI response success user_mail=%s conversation_id=%s len=%d preview=\"%s\"",
                user_info.user_mail,
                user_info.conversation_id,
                len(response or ""),
                response_preview or "<empty>",
            )

            await turn_context.send_activity(
                Activity(
                    type=ActivityTypes.message,
                    text=response,
                    suggested_actions=(
                        SuggestedActions(actions=suggested_actions)
                        if suggested_actions
                        else None
                    ),
                )
            )
        except OpenAIServiceError as openai_error:
            error_hint = str(openai_error) or "OpenAI 服務暫時無法回應"
            fallback_text = (
                "⚠️ 呼叫 OpenAI 模型時發生問題，暫時無法提供回覆。\n"
                "請稍後再試，或輸入 @model 改用其他模型。"
            )
            await turn_context.send_activity(
                Activity(type=ActivityTypes.message, text=fallback_text)
            )
            self.logger.warning(
                "OpenAIServiceError while responding user_mail=%s conversation_id=%s error=%s",
                user_info.user_mail,
                user_info.conversation_id,
                error_hint,
            )
        except Exception as unexpected_error:
            await turn_context.send_activity(
                Activity(
                    type=ActivityTypes.message,
                    text="⚠️ 目前無法取得模型回覆，請稍後再試一次。",
                )
            )
            self.logger.exception(
                "Unexpected error in _handle_direct_openai_response user_mail=%s conversation_id=%s",
                user_info.user_mail,
                user_info.conversation_id,
            )

    async def _try_attach_images(self, turn_context: TurnContext, user_info: BotInteractionDTO) -> bool:
        """If message has file attachments and user has a recent IT task, upload them to Asana.
        Returns True if handled (uploaded or error responded), else False.
        """
        try:
            atts = turn_context.activity.attachments or []
            files = []
            # Accept general file attachments
            for a in atts:
                ctype = (getattr(a, "content_type", None) or "").lower()
                name = getattr(a, "name", None) or ""
                url = getattr(a, "content_url", None) or getattr(a, "contentUrl", None)
                # Teams file info attachments (e.g., application/vnd.microsoft.teams.file.download.info)
                if not url and ctype == "application/vnd.microsoft.teams.file.download.info":
                    content = getattr(a, "content", None) or {}
                    if isinstance(content, dict):
                        url = content.get("downloadUrl")
                        # Try resolve a better filename
                        if not name:
                            name = content.get("name") or name
                        ft = content.get("fileType")
                        generic = (not name) or (name.lower() in ("file", "file.bin", "image", "image.jpg", "original", "upload", "upload.bin")) or ("." not in name)
                        if generic and ft:
                            from datetime import datetime
                            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
                            name = f"upload_{ts}.{ft}"
                # Skip card attachments
                if ctype.startswith("application/vnd.microsoft.card"):
                    continue
                # Handle data URL (e.g., Bot Emulator inline data)
                if url and str(url).startswith("data:"):
                    try:
                        header, b64data = str(url).split(",", 1)
                        import base64
                        mime = "application/octet-stream"
                        if ":" in header and ";" in header:
                            mime = header.split(":", 1)[1].split(";", 1)[0] or mime
                        # If mime still unknown, try to sniff from bytes
                        data_bytes = base64.b64decode(b64data)
                        if not mime or mime == "application/octet-stream":
                            try:
                                # Simple magic sniff for common types
                                if data_bytes.startswith(b"\x89PNG\r\n\x1a\n"):
                                    mime = "image/png"
                                elif data_bytes.startswith(b"\xff\xd8"):
                                    mime = "image/jpeg"
                                elif data_bytes.startswith(b"GIF8"):
                                    mime = "image/gif"
                                elif data_bytes.startswith(b"RIFF") and b"WEBP" in data_bytes[:16]:
                                    mime = "image/webp"
                                elif data_bytes.startswith(b"%PDF"):
                                    mime = "application/pdf"
                                elif data_bytes.startswith(b"PK\x03\x04"):
                                    mime = "application/zip"
                            except Exception:
                                pass
                        # Infer extension if name missing or generic (e.g., 'original')
                        generic = (not name) or (name.lower() in ("file", "file.bin", "image", "image.jpg", "original", "upload", "upload.bin")) or ("." not in name)
                        if generic:
                            ext_map = {
                                "image/png": "png",
                                "image/jpeg": "jpg",
                                "image/jpg": "jpg",
                                "image/gif": "gif",
                                "image/webp": "webp",
                                "image/bmp": "bmp",
                                "image/heic": "heic",
                                "application/pdf": "pdf",
                                "application/zip": "zip",
                                "text/plain": "txt",
                                "application/msword": "doc",
                                "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
                                "application/vnd.ms-excel": "xls",
                                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
                                "application/vnd.ms-powerpoint": "ppt",
                                "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
                            }
                            ext = ext_map.get(mime, "bin")
                            from datetime import datetime
                            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
                            name = f"screenshot_{ts}.{ext}"
                        files.append({"data": data_bytes, "name": name or "file.bin", "ctype": mime})
                        continue
                    except Exception:
                        pass
                if url:
                    # Derive filename from URL if name missing or generic
                    generic = (not name) or (name.lower() in ("file", "file.bin", "image", "image.jpg", "original", "upload", "upload.bin")) or ("." not in name)
                    if generic:
                        try:
                            from urllib.parse import urlparse, unquote
                            path = urlparse(url).path
                            base = path.rsplit("/", 1)[-1]
                            if base:
                                name = unquote(base)
                        except Exception:
                            pass
                    # If still generic after URL parse, synthesize from content-type
                    generic2 = (not name) or (name.lower() in ("file", "file.bin", "image", "image.jpg", "original", "upload", "upload.bin")) or ("." not in name)
                    if generic2:
                        ext_map = {
                            "image/png": "png",
                            "image/jpeg": "jpg",
                            "image/jpg": "jpg",
                            "image/gif": "gif",
                            "image/webp": "webp",
                            "image/bmp": "bmp",
                            "image/heic": "heic",
                            "application/pdf": "pdf",
                            "application/zip": "zip",
                            "text/plain": "txt",
                            "application/msword": "doc",
                            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
                            "application/vnd.ms-excel": "xls",
                            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
                            "application/vnd.ms-powerpoint": "ppt",
                            "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
                        }
                        ext = ext_map.get(ctype, "bin")
                        from datetime import datetime
                        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
                        name = f"upload_{ts}.{ext}"
                    files.append({"url": url, "name": name or "file.bin", "ctype": ctype or "application/octet-stream"})
            if not files:
                return False

            from core.container import get_container
            from features.it_support.service import ITSupportService
            svc: ITSupportService = get_container().get(ITSupportService)
            gid = svc.get_recent_task_gid(user_info.user_mail)
            if not gid:
                # 沒有最近的 IT 單，不攔截，讓後續流程（如圖片解析）處理
                return False

            # Try attach each image by URL
            ok = 0
            for a in files:
                name = a.get("name") or "file.bin"
                ctype = a.get("ctype") or "application/octet-stream"
                if "data" in a:
                    result = await svc.attach_image_bytes(user_info.user_mail, a["data"], name, ctype)
                else:
                    url = a["url"]
                    result = await svc.attach_image_from_url(user_info.user_mail, url, name, ctype)
                if result.get("success"):
                    ok += 1
                else:
                    await turn_context.send_activity(
                        Activity(type=ActivityTypes.message, text=f"❌ 檔案上傳失敗：{result.get('error')}")
                    )

            if ok:
                await turn_context.send_activity(
                    Activity(type=ActivityTypes.message, text=f"✅ 已上傳 {ok} 個檔案至最近的 IT 單")
                )
            return True
        except Exception as e:
            await turn_context.send_activity(
                Activity(type=ActivityTypes.message, text=f"❌ 上傳檔案時發生錯誤：{str(e)}")
            )
            return True

    def _parse_attachment_files(self, turn_context: TurnContext) -> List[dict]:
        """從 activity.attachments 解析出檔案清單。
        複用 _try_attach_images 已驗證的解析邏輯。
        Returns list of dict: {'data': bytes, 'url': str, 'name': str, 'ctype': str}
        """
        import base64
        atts = turn_context.activity.attachments or []
        files = []
        for a in atts:
            ctype = (getattr(a, "content_type", None) or "").lower()
            name = getattr(a, "name", None) or ""
            url = getattr(a, "content_url", None) or getattr(a, "contentUrl", None)

            # Teams file info attachments
            if ctype == "application/vnd.microsoft.teams.file.download.info":
                content = getattr(a, "content", None) or {}
                if isinstance(content, dict):
                    url = content.get("downloadUrl") or url
                    if not name:
                        name = content.get("name") or name
                    ft = content.get("fileType")
                    if ft and not name:
                        from datetime import datetime
                        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
                        name = f"upload_{ts}.{ft}"
                    # 根據 fileType 或檔名推斷正確的 MIME type
                    ext = (ft or "").lower() if ft else ""
                    if not ext and name:
                        ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
                    mime_map = {
                        "pdf": "application/pdf",
                        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        "doc": "application/msword",
                        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        "xls": "application/vnd.ms-excel",
                        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        "ppt": "application/vnd.ms-powerpoint",
                        "png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg",
                        "gif": "image/gif", "bmp": "image/bmp", "webp": "image/webp",
                        "txt": "text/plain", "csv": "text/csv", "json": "application/json",
                        "xml": "text/xml", "md": "text/markdown", "log": "text/plain",
                    }
                    if ext in mime_map:
                        ctype = mime_map[ext]

            # Skip card attachments and HTML previews
            if ctype.startswith("application/vnd.microsoft.card"):
                continue
            if ctype == "text/html":
                continue

            # data URL (Bot Emulator)
            if url and str(url).startswith("data:"):
                try:
                    header, b64data = str(url).split(",", 1)
                    mime = "application/octet-stream"
                    if ":" in header and ";" in header:
                        mime = header.split(":", 1)[1].split(";", 1)[0] or mime
                    data_bytes = base64.b64decode(b64data)
                    if not mime or mime == "application/octet-stream":
                        if data_bytes.startswith(b"\x89PNG\r\n\x1a\n"): mime = "image/png"
                        elif data_bytes.startswith(b"\xff\xd8"): mime = "image/jpeg"
                        elif data_bytes.startswith(b"GIF8"): mime = "image/gif"
                        elif data_bytes.startswith(b"BM"): mime = "image/bmp"
                        elif data_bytes.startswith(b"%PDF"): mime = "application/pdf"
                    files.append({"data": data_bytes, "name": name or "file", "ctype": mime})
                    continue
                except Exception:
                    continue

            # HTTP URL — 記錄 URL，稍後下載
            if url and str(url).startswith("http"):
                files.append({"url": url, "name": name or "file", "ctype": ctype or "application/octet-stream"})

        return files

    async def _download_parsed_files(self, files: List[dict]) -> List[dict]:
        """下載 _parse_attachment_files 回傳的檔案清單，使用 ITSupportService 已驗證的下載方法。
        Returns list of dict: {'bytes': bytes, 'mime_type': str, 'name': str}
        """
        import httpx
        results = []
        for f in files:
            name = f.get("name") or "file"
            ctype = f.get("ctype") or "application/octet-stream"

            # 已有 bytes（data URL）
            if "data" in f:
                results.append({"bytes": f["data"], "mime_type": ctype, "name": name})
                continue

            url = f.get("url")
            if not url:
                continue

            try:
                # 使用 ITSupportService 的 Bot Framework token 認證（與 IT 附件上傳相同的方法）
                headers = {}
                from core.container import get_container
                from features.it_support.service import ITSupportService
                svc: ITSupportService = get_container().get(ITSupportService)

                if svc._is_botframework_protected_url(url):
                    token = await svc._get_botframework_token()
                    if token:
                        headers["Authorization"] = f"Bearer {token}"
                        headers["Accept"] = "*/*"
                        self.logger.info("Using BotFramework token for attachment download")

                async with httpx.AsyncClient(timeout=30.0, follow_redirects=True) as client:
                    resp = await client.get(url, headers=headers)
                    self.logger.info(
                        "Attachment download status=%d content_type=%s size=%d url=%s",
                        resp.status_code, resp.headers.get("content-type", ""), len(resp.content), url[:80],
                    )
                    if resp.status_code == 200 and resp.content:
                        data = resp.content
                        # 用 response header 確定 mime type
                        resp_ct = (resp.headers.get("content-type") or "").split(";")[0].strip().lower()
                        if resp_ct and resp_ct != "application/octet-stream" and "*" not in resp_ct:
                            ctype = resp_ct
                        # magic byte sniffing for uncertain types
                        if not ctype or ctype == "application/octet-stream" or "*" in ctype:
                            if data.startswith(b"\x89PNG\r\n\x1a\n"): ctype = "image/png"
                            elif data.startswith(b"\xff\xd8"): ctype = "image/jpeg"
                            elif data.startswith(b"GIF8"): ctype = "image/gif"
                            elif data.startswith(b"RIFF") and b"WEBP" in data[:16]: ctype = "image/webp"
                            elif data.startswith(b"BM"): ctype = "image/bmp"
                            elif data.startswith(b"%PDF"): ctype = "application/pdf"
                        results.append({"bytes": data, "mime_type": ctype, "name": name})
                    else:
                        self.logger.warning("Attachment download failed status=%d url=%s", resp.status_code, url[:100])
            except Exception as e:
                self.logger.warning("Failed to download attachment: %s url=%s", e, url[:100])

        return results

    _ATTACHMENT_MAX_SIZE_MB = 20  # 單一附件大小上限

    async def _handle_attachment_analysis(
        self, turn_context: TurnContext, user_info: BotInteractionDTO,
        attachments: List[dict],
    ) -> None:
        """用 AI 解析使用者上傳的附件（圖片用 Vision，文件用文字擷取）。支援多筆附件。"""
        import base64

        # 檔案大小檢查：超過上限的檔案給予友善提示
        max_bytes = int(self._ATTACHMENT_MAX_SIZE_MB * 1024 * 1024)
        oversized = [a for a in attachments if len(a.get("bytes", b"")) > max_bytes]
        attachments = [a for a in attachments if len(a.get("bytes", b"")) <= max_bytes]

        if oversized:
            names = ", ".join(a["name"] for a in oversized)
            sizes = ", ".join(f"{len(a['bytes'])/1024/1024:.1f}MB" for a in oversized)
            await turn_context.send_activity(Activity(
                type=ActivityTypes.message,
                text=f"⚠️ 檔案過大無法解析：{names}（{sizes}）。\n"
                     f"目前支援的檔案大小上限為 {self._ATTACHMENT_MAX_SIZE_MB}MB。\n"
                     f"💡 建議：可將大型 PDF 拆分為較小的檔案後重新上傳，或僅上傳需要解析的部分頁面。",
            ))
            if not attachments:
                return

        user_text = (user_info.message_text or "").strip()
        openai_client = self.conversation_service.openai_client
        vision_model = self.config.openai.vision_model

        # 分離圖片與文件
        images = [a for a in attachments if a["mime_type"].startswith("image/")]
        files = [a for a in attachments if not a["mime_type"].startswith("image/")]

        try:
            # --- 有圖片：用 Vision 一次送所有圖片 ---
            if images:
                prompt = user_text if user_text else "請描述這些圖片的內容。如果是錯誤截圖，請說明可能的問題和建議解法。"
                content_parts = [{"type": "text", "text": prompt}]
                for img in images:
                    b64 = base64.b64encode(img["bytes"]).decode("utf-8")
                    content_parts.append({
                        "type": "image_url",
                        "image_url": {"url": f"data:{img['mime_type']};base64,{b64}"},
                    })

                messages = [{"role": "user", "content": content_parts}]
                response = await openai_client.chat_completion(
                    messages=messages, model=vision_model, max_tokens=1500,
                )
                await turn_context.send_activity(
                    Activity(type=ActivityTypes.message, text=response)
                )

            # --- 有文件：逐一擷取文字，合併後送 AI ---
            if files:
                all_texts = []
                scanned_pdfs = []  # 掃描型 PDF，需用 Vision 處理
                unsupported = []
                for f in files:
                    # 判斷是否為 PDF（掃描型會回傳 None）
                    is_pdf = (f["mime_type"] == "application/pdf" or f["name"].lower().endswith(".pdf"))
                    file_text = self._extract_text_from_file(f["bytes"], f["mime_type"], f["name"])
                    if file_text:
                        all_texts.append(f"📄 **{f['name']}**\n{file_text}")
                    elif is_pdf:
                        scanned_pdfs.append(f)
                    else:
                        unsupported.append(f["name"])

                if unsupported:
                    await turn_context.send_activity(Activity(
                        type=ActivityTypes.message,
                        text=f"⚠️ 無法解析：{', '.join(unsupported)}。目前支援圖片（PNG/JPG/BMP/GIF/WebP）、PDF、Word(.docx)、Excel(.xlsx/.xls)、PowerPoint(.pptx)、純文字檔。",
                    ))

                if all_texts:
                    combined = "\n\n---\n\n".join(all_texts)
                    max_chars = 12000
                    if len(combined) > max_chars:
                        combined = combined[:max_chars] + "\n\n...（內容過長，已截斷）"
                    prompt = user_text if user_text else "以下是上傳檔案的內容，請分析並摘要重點："
                    messages = [{"role": "user", "content": f"{prompt}\n\n---\n{combined}"}]
                    response = await openai_client.chat_completion(
                        messages=messages, max_tokens=1500,
                    )
                    await turn_context.send_activity(
                        Activity(type=ActivityTypes.message, text=response)
                    )

                # --- 掃描型 PDF：分批轉圖片 → Vision 辨識 → 合併摘要 ---
                for pdf_file in scanned_pdfs:
                    pdf_images, total_pages = self._pdf_to_images(pdf_file["bytes"], pdf_file["name"])
                    if not pdf_images:
                        await turn_context.send_activity(Activity(
                            type=ActivityTypes.message,
                            text=f"⚠️ 無法解析掃描型 PDF：{pdf_file['name']}",
                        ))
                        continue

                    # 分批送 Vision（每批 _VISION_CHUNK_SIZE 頁）
                    chunk_size = self._VISION_CHUNK_SIZE
                    chunks = [pdf_images[i:i + chunk_size] for i in range(0, len(pdf_images), chunk_size)]
                    chunk_texts = []

                    for idx, chunk in enumerate(chunks):
                        page_start = idx * chunk_size + 1
                        page_end = page_start + len(chunk) - 1
                        ocr_prompt = (
                            f"以下是 PDF「{pdf_file['name']}」第 {page_start}~{page_end} 頁的圖片，"
                            f"請完整辨識其中的所有文字內容，保留原始結構。"
                        )
                        content_parts: list = [{"type": "text", "text": ocr_prompt}]
                        for img_b64 in chunk:
                            content_parts.append({
                                "type": "image_url",
                                "image_url": {"url": f"data:image/png;base64,{img_b64}"},
                            })

                        messages = [{"role": "user", "content": content_parts}]
                        chunk_result = await openai_client.chat_completion(
                            messages=messages, model=vision_model, max_tokens=4000,
                        )
                        chunk_texts.append(f"【第 {page_start}~{page_end} 頁】\n{chunk_result}")

                    # 合併所有 chunk 結果
                    combined_ocr = "\n\n---\n\n".join(chunk_texts)
                    truncate_note = ""
                    if len(pdf_images) < total_pages:
                        truncate_note = f"\n\n...（僅辨識前 {len(pdf_images)} 頁，共 {total_pages} 頁）"

                    # 如果只有一個 chunk，直接回傳辨識結果
                    if len(chunks) == 1:
                        await turn_context.send_activity(Activity(
                            type=ActivityTypes.message,
                            text=chunk_texts[0] + truncate_note,
                        ))
                    else:
                        # 多個 chunk → 再送一次 AI 做總結摘要
                        summary_prompt = user_text if user_text else (
                            f"以下是 PDF「{pdf_file['name']}」經 OCR 辨識後的文字內容（共 {len(pdf_images)} 頁），"
                            f"請分析並摘要重點。"
                        )
                        summary_messages = [{"role": "user", "content": f"{summary_prompt}\n\n---\n{combined_ocr}{truncate_note}"}]
                        summary = await openai_client.chat_completion(
                            messages=summary_messages, max_tokens=4000,
                        )
                        await turn_context.send_activity(
                            Activity(type=ActivityTypes.message, text=summary)
                        )

            # 附件分析完成後顯示支援格式提示
            tip = ("💡 支援格式：圖片（PNG/JPG/BMP/GIF/WebP）、PDF、Word(.docx)、"
                   "Excel(.xlsx/.xls)、PowerPoint(.pptx)、純文字（TXT/CSV/JSON/XML/MD）")
            await turn_context.send_activity(Activity(type=ActivityTypes.message, text=tip))

        except Exception:
            names = ", ".join(a["name"] for a in attachments)
            self.logger.exception("附件解析失敗 user_mail=%s files=%s", user_info.user_mail, names)
            await turn_context.send_activity(Activity(
                type=ActivityTypes.message, text="❌ 檔案解析失敗，請稍後再試。",
            ))

    # 不支援的舊格式，提示使用者另存新版
    _LEGACY_FORMATS = {
        ".doc": ".docx", ".xls": ".xlsx", ".ppt": ".pptx",
    }

    def _extract_text_from_file(self, data: bytes, mime: str, name: str) -> Optional[str]:
        """從檔案 bytes 中擷取文字內容"""
        lower_name = name.lower()

        # 舊格式提示另存新版
        for old_ext, new_ext in self._LEGACY_FORMATS.items():
            if lower_name.endswith(old_ext) and not lower_name.endswith(new_ext):
                return f"⚠️ 不支援舊版 {old_ext} 格式，請另存為 {new_ext} 後重新上傳。"

        try:
            # PDF
            if mime == "application/pdf" or lower_name.endswith(".pdf"):
                from PyPDF2 import PdfReader
                import io
                reader = PdfReader(io.BytesIO(data))
                pages = [p.extract_text() or "" for p in reader.pages[:20]]
                text = "\n".join(pages).strip()
                if text:
                    return f"📄 PDF格式: 文字\n\n{text}"
                # 文字擷取失敗 → 標記為掃描型 PDF，由 _handle_attachment_analysis 用 Vision 處理
                return None

            # Word (.docx)
            if "wordprocessingml" in mime or lower_name.endswith(".docx"):
                from docx import Document
                import io
                doc = Document(io.BytesIO(data))
                return "\n".join(p.text for p in doc.paragraphs).strip() or None

            # Excel (.xlsx)
            if "spreadsheetml" in mime or lower_name.endswith(".xlsx"):
                import openpyxl
                import io
                wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
                lines = []
                for ws in wb.worksheets[:3]:  # 最多讀 3 個 sheet
                    lines.append(f"[Sheet: {ws.title}]")
                    for row in ws.iter_rows(max_row=100, values_only=True):
                        cells = [str(c) if c is not None else "" for c in row]
                        lines.append("\t".join(cells))
                return "\n".join(lines).strip() or None

            # Excel (.xls 舊格式)
            if "ms-excel" in mime or lower_name.endswith(".xls"):
                import xlrd
                import io
                wb = xlrd.open_workbook(file_contents=data)
                lines = []
                for ws in wb.sheets()[:3]:
                    lines.append(f"[Sheet: {ws.name}]")
                    for row_idx in range(min(ws.nrows, 100)):
                        cells = [str(ws.cell_value(row_idx, col)) for col in range(ws.ncols)]
                        lines.append("\t".join(cells))
                return "\n".join(lines).strip() or None

            # PowerPoint (.pptx)
            if "presentationml" in mime or lower_name.endswith(".pptx"):
                from pptx import Presentation
                import io
                prs = Presentation(io.BytesIO(data))
                lines = []
                for i, slide in enumerate(prs.slides[:50], 1):
                    texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                t = paragraph.text.strip()
                                if t:
                                    texts.append(t)
                    if texts:
                        lines.append(f"[Slide {i}]")
                        lines.extend(texts)
                return "\n".join(lines).strip() or None

            # 純文字（含 .md）
            if (mime.startswith("text/") or "json" in mime or "xml" in mime
                    or lower_name.endswith((".txt", ".csv", ".log", ".json", ".xml", ".md"))):
                return data.decode("utf-8", errors="replace").strip() or None

        except Exception as e:
            self.logger.warning("擷取檔案文字失敗 name=%s: %s", name, e)

        return None

    _VISION_CHUNK_SIZE = 5  # 每批送 Vision 的頁數
    _VISION_PDF_DPI = 150  # PDF 轉圖片的 DPI

    def _pdf_to_images(self, data: bytes, name: str) -> tuple:
        """將掃描型 PDF 轉換為 base64 PNG 圖片清單，供 Vision 辨識。

        Returns:
            (images: list[str], total_pages: int) — images 為 base64 PNG 字串清單
        """
        import base64
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(stream=data, filetype="pdf")
            total_pages = len(doc)
            max_pages = total_pages
            images = []

            for i in range(max_pages):
                page = doc[i]
                pix = page.get_pixmap(dpi=self._VISION_PDF_DPI)
                img_bytes = pix.tobytes("png")
                b64 = base64.b64encode(img_bytes).decode("utf-8")
                images.append(b64)

            doc.close()

            self.logger.info("PDF 轉圖片完成: name=%s, pages=%d", name, len(images))
            return images, total_pages
        except Exception as e:
            self.logger.warning("PDF 轉圖片失敗 name=%s: %s", name, e)
            return [], 0

    async def _send_welcome_message(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """發送歡迎訊息"""
        language = determine_language(user_info.user_mail)
        welcome_msg = {
            "zh": "🎉 歡迎使用台灣林內 GPT！\n我可以協助您管理待辦事項、預約會議室等功能。",
            "en": "🎉 Welcome to Taiwan Rinnai GPT!\nI can help you manage todos, book meeting rooms, and more.",
            "ja": "🎉 台湾リンナイGPTへようこそ！\nタスク管理、会議室予約などをサポートします。",
        }.get(language, "🎉 歡迎使用台灣林內 GPT！")

        include_model = not self.config.openai.use_azure
        help_card = self.help_card_builder.build_help_card(language, welcome_msg, include_model_option=include_model)
        await turn_context.send_activity(help_card)

    async def _show_add_todo_card(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """顯示新增待辦事項卡片"""
        language = determine_language(user_info.user_mail)
        card = self.todo_card_builder.build_add_todo_card(language)
        await turn_context.send_activity(card)

    async def _show_todo_list(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """顯示待辦事項清單"""
        todos = await self.todo_service.get_user_todos(
            user_info.user_mail, include_completed=False
        )
        language = determine_language(user_info.user_mail)

        if todos:
            card = self.todo_card_builder.build_todo_list_card(todos, language)
            await turn_context.send_activity(card)
            # 小提示（對齊 app_bak 互動風格）
            hint_msg = (
                "💡 小提示：下次可以直接輸入 `@ls` 快速查看待辦清單"
                if language == "zh-TW"
                else "💡 ヒント：次回は `@ls` で素早くTODOリストを確認できます"
            )
            await turn_context.send_activity(
                Activity(type=ActivityTypes.message, text=hint_msg)
            )
        else:
            suggested_actions = get_suggested_replies("無待辦事項", user_info.user_mail)
            await turn_context.send_activity(
                Activity(
                    type=ActivityTypes.message,
                    text="🎉 目前沒有待辦事項",
                    suggested_actions=(
                        SuggestedActions(actions=suggested_actions)
                        if suggested_actions
                        else None
                    ),
                )
            )

    async def _show_room_booking_options(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """顯示會議室預約選項"""
        language = determine_language(user_info.user_mail)
        card = self.meeting_card_builder.build_room_booking_card(language)
        await turn_context.send_activity(card)
        # 小提示
        hint_msg = (
            "💡 小提示：也可以使用 `@book-room` 快速開啟預約表單"
            if language == "zh-TW"
            else "💡 ヒント：`@book-room` でも素早く予約フォームを開けます"
        )
        await turn_context.send_activity(
            Activity(type=ActivityTypes.message, text=hint_msg)
        )

    async def _show_my_bookings(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """顯示我的預約"""
        bookings = await self.meeting_service.get_user_meetings(user_info.user_mail)
        language = determine_language(user_info.user_mail)

        if bookings:
            card = self.meeting_card_builder.build_my_bookings_card(bookings, language)
            await turn_context.send_activity(card)
            hint_msg = (
                "💡 小提示：也可以使用 `@check-booking` 快速查看預約"
                if language == "zh-TW"
                else "💡 ヒント：`@check-booking` でも素早く予約を確認できます"
            )
            await turn_context.send_activity(
                Activity(type=ActivityTypes.message, text=hint_msg)
            )
        else:
            await turn_context.send_activity(
                Activity(
                    type=ActivityTypes.message,
                    text="📅 目前沒有預約的會議室",
                )
            )

    async def _show_cancel_booking_options(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """顯示取消預約選項"""
        bookings = await self.meeting_service.get_user_meetings(user_info.user_mail)
        language = determine_language(user_info.user_mail)

        if bookings:
            card = self.meeting_card_builder.build_cancel_booking_card(
                bookings, language
            )
            await turn_context.send_activity(card)
            hint_msg = (
                "💡 小提示：也可以使用 `@cancel-booking` 快速取消預約"
                if language == "zh-TW"
                else "💡 ヒント：`@cancel-booking` でも素早く予約をキャンセルできます"
            )
            await turn_context.send_activity(
                Activity(type=ActivityTypes.message, text=hint_msg)
            )
        else:
            await turn_context.send_activity(
                Activity(
                    type=ActivityTypes.message,
                    text="📅 目前沒有可取消的會議室預約",
                )
            )

    async def _show_user_info(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """顯示用戶資訊（擴充為更完整的個人檔案）。"""
        # 基本欄位
        lines = [
            "👤 **用戶資訊**",
            f"• 姓名: {user_info.user_name or '未知'}",
            f"• 郵箱: {user_info.user_mail}",
            f"• 用戶ID: {user_info.user_id}",
        ]

        # 系統狀態（模式/模型/意圖分析）
        if self.config.openai.use_azure:
            mode_text = "Azure OpenAI"
            model_text = "o1-mini（固定）"
        else:
            mode_text = "OpenAI"
            try:
                from app import user_model_preferences
                model_text = user_model_preferences.get(user_info.user_mail, self.config.openai.model)
            except Exception:
                model_text = self.config.openai.model
        lines += [
            "",
            "🤖 **系統狀態**",
            f"• 模式: {mode_text}",
            f"• 模型: {model_text}",
            f"• AI 意圖分析: {'啟用' if self.config.enable_ai_intent_analysis else '停用'}",
        ]

        # 透過 Graph 取得更完整的資料（部門/職稱/電話/分機等）
        try:
            from core.container import get_container
            from infrastructure.external.graph_api_client import GraphAPIClient

            container = get_container()
            graph: GraphAPIClient = container.get(GraphAPIClient)

            async with graph:
                profile = await graph.get_user_rich_profile(user_info.user_mail)
                manager = await graph.get_user_manager(user_info.user_mail)
                extension = graph.try_extract_extension(profile)

            dept = profile.get("department")
            title = profile.get("jobTitle")
            company = profile.get("companyName")
            office = profile.get("officeLocation")
            mobile = profile.get("mobilePhone")
            biz_phones = profile.get("businessPhones") or []
            employee_id = profile.get("employeeId")
            org = profile.get("employeeOrgData") or {}
            cost_center = org.get("costCenter")
            division = org.get("division")
            mgr_name = ((manager or {}).get("displayName")) if manager else None

            if dept:
                lines.append(f"• 部門: {dept}")
            if title:
                lines.append(f"• 職稱: {title}")
            if company:
                lines.append(f"• 公司: {company}")
            if office:
                lines.append(f"• 辦公室: {office}")
            if extension:
                lines.append(f"• 分機: {extension}")
            if mobile:
                lines.append(f"• 手機: {mobile}")
            if biz_phones:
                lines.append(f"• 電話: {'、'.join(biz_phones)}")
            if employee_id:
                lines.append(f"• 員工編號: {employee_id}")
            if cost_center or division:
                cc = f"成本中心: {cost_center}" if cost_center else None
                dv = f"事業群: {division}" if division else None
                org_line = "、".join([v for v in [cc, dv] if v])
                if org_line:
                    lines.append(f"• {org_line}")
            if mgr_name:
                lines.append(f"• 主管: {mgr_name}")
        except Exception as e:
            # 若 Graph 取資料失敗，僅顯示基本欄位
            pass

        await turn_context.send_activity(
            Activity(type=ActivityTypes.message, text="\n".join(lines))
        )

    async def _show_bot_intro(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """顯示機器人介紹"""
        language = determine_language(user_info.user_mail)
        intro_card = self.help_card_builder.build_bot_intro_card(language)
        await turn_context.send_activity(intro_card)

    async def _show_model_selection(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """顯示模型選擇"""
        if self.config.openai.use_azure:
            await turn_context.send_activity(
                Activity(
                    type=ActivityTypes.message,
                    text="ℹ️ 目前使用 Azure OpenAI 服務\n📱 模型：o1-mini（固定）\n⚡ 此模式不支援模型切換",
                )
            )
            return

        # OpenAI 模式：顯示模型選擇卡片
        card = self.model_card_builder.build_model_selection_card(user_info.user_mail)
        await turn_context.send_activity(card)

    async def _handle_add_todo_card(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """處理新增待辦事項卡片"""
        content = turn_context.activity.value.get("todoContent", "").strip()
        if content:
            await self._handle_smart_todo_add(turn_context, user_info, content)

    async def _handle_smart_todo_add(
        self, turn_context: TurnContext, user_info: BotInteractionDTO, content: str
    ) -> None:
        """智能新增待辦事項"""
        try:
            todo, similar_todos = await self.todo_service.smart_create_todo(
                user_info.user_mail, content
            )

            if similar_todos:
                # 有相似的待辦事項，顯示確認卡片
                language = determine_language(user_info.user_mail)
                card = self.todo_card_builder.build_similar_todos_confirmation_card(
                    content, similar_todos, language
                )
                await turn_context.send_activity(card)
            elif todo:
                # 成功新增
                await turn_context.send_activity(
                    Activity(
                        type=ActivityTypes.message,
                        text=f"✅ 已新增待辦事項：{todo.content}",
                    )
                )
                # 小提示
                language = determine_language(user_info.user_mail)
                hint_msg = (
                    "💡 小提示：下次可以使用 `@add 內容` 快速新增待辦"
                    if language == "zh-TW"
                    else "💡 ヒント：次回は `@add 内容` で素早くTODOを追加できます"
                )
                await turn_context.send_activity(
                    Activity(type=ActivityTypes.message, text=hint_msg)
                )
        except Exception as e:
            await turn_context.send_activity(
                Activity(
                    type=ActivityTypes.message, text=f"❌ 新增待辦事項失敗：{str(e)}"
                )
            )

    async def _handle_complete_todos(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """處理完成待辦事項"""
        completed_indices_str = turn_context.activity.value.get("completedTodos", "")

        if not completed_indices_str:
            return

        try:
            # 解析完成的待辦事項索引
            completed_indices = [
                int(idx.strip())
                for idx in completed_indices_str.split(",")
                if idx.strip()
            ]

            if completed_indices:
                completed_todos = await self.todo_service.batch_complete_todos(
                    completed_indices, user_info.user_mail
                )

                if completed_todos:
                    completed_text = "\n".join(
                        [f"• {todo.content}" for todo in completed_todos]
                    )
                    await turn_context.send_activity(
                        Activity(
                            type=ActivityTypes.message,
                            text=f"✅ 已完成 {len(completed_todos)} 項待辦事項：\n{completed_text}",
                        )
                    )
                else:
                    await turn_context.send_activity(
                        Activity(
                            type=ActivityTypes.message, text="❌ 無法完成指定的待辦事項"
                        )
                    )
        except Exception as e:
            await turn_context.send_activity(
                Activity(
                    type=ActivityTypes.message,
                    text=f"❌ 完成待辦事項時發生錯誤：{str(e)}",
                )
            )

    async def _handle_book_room(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """處理會議室預約"""
        # 從卡片提取預約信息
        booking_data = {
            "room_id": turn_context.activity.value.get("selectedRoom"),
            "date": turn_context.activity.value.get("selectedDate"),
            "start_time": turn_context.activity.value.get("startTime"),
            "end_time": turn_context.activity.value.get("endTime"),
            "subject": turn_context.activity.value.get("subject", ""),
            # "attendees": turn_context.activity.value.get("attendees", "")
        }

        try:
            booking_result = await self.meeting_service.book_meeting_room(
                user_info.user_mail, booking_data
            )

            if booking_result.get("success"):
                await turn_context.send_activity(
                    Activity(
                        type=ActivityTypes.message,
                        text=f"✅ 成功預約會議室：{booking_data['subject']}",
                    )
                )
            else:
                error_msg = booking_result.get("error", "未知錯誤")
                await turn_context.send_activity(
                    Activity(
                        type=ActivityTypes.message, text=f"❌ 預約失敗：{error_msg}"
                    )
                )
        except Exception as e:
            await turn_context.send_activity(
                Activity(
                    type=ActivityTypes.message,
                    text=f"❌ 預約會議室時發生錯誤：{str(e)}",
                )
            )

    async def _handle_cancel_booking(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """處理取消預約"""
        booking_id = turn_context.activity.value.get("selectedBooking")

        if not booking_id:
            return

        try:
            cancel_result = await self.meeting_service.cancel_meeting(
                user_info.user_mail, booking_id
            )

            if cancel_result.get("success"):
                await turn_context.send_activity(
                    Activity(type=ActivityTypes.message, text="✅ 已成功取消會議室預約")
                )
            else:
                error_msg = cancel_result.get("error", "未知錯誤")
                await turn_context.send_activity(
                    Activity(
                        type=ActivityTypes.message, text=f"❌ 取消預約失敗：{error_msg}"
                    )
                )
        except Exception as e:
            await turn_context.send_activity(
                Activity(
                    type=ActivityTypes.message, text=f"❌ 取消預約時發生錯誤：{str(e)}"
                )
            )

    async def _handle_model_selection(
        self, turn_context: TurnContext, user_info: BotInteractionDTO
    ) -> None:
        """處理模型選擇"""
        selected_model = turn_context.activity.value.get("selectedModel")

        if selected_model:
            # 更新用戶模型偏好
            # 這裡應該透過用戶服務來更新
            from app import user_model_preferences, MODEL_INFO

            user_model_preferences[user_info.user_mail] = selected_model

            model_info = MODEL_INFO.get(selected_model, {})
            await turn_context.send_activity(
                Activity(
                    type=ActivityTypes.message,
                    text=f"✅ 已切換到模型：{selected_model}\n{model_info.get('use_case', '')}",
                )
            )

    async def _send_proactive_message(self, conversation_ref, text: str) -> None:
        """透過 proactive message 發送文字訊息"""
        import os
        from core.container import get_container
        from infrastructure.bot.bot_adapter import CustomBotAdapter

        bot_adapter: CustomBotAdapter = get_container().get(CustomBotAdapter)
        bot_app_id = os.getenv("BOT_APP_ID") or os.getenv("MICROSOFT_APP_ID") or ""

        async def send(turn_context):
            await turn_context.send_activity(
                Activity(type=ActivityTypes.message, text=text)
            )

        await bot_adapter.adapter.continue_conversation(
            conversation_ref, send, bot_app_id
        )

    async def _send_error_response(self, turn_context: TurnContext) -> None:
        """發送錯誤回應"""
        await turn_context.send_activity(
            Activity(
                type=ActivityTypes.message,
                text="❌ 處理您的請求時發生錯誤，請稍後再試。",
            )
        )

    async def show_help_options(
        self, turn_context: TurnContext, welcome_msg: str = None
    ):
        """顯示幫助選項（從原始 app_bak.py 遷移）"""
        # 取得用戶語系
        user_id = turn_context.activity.from_property.id
        user_mail = await get_user_email(turn_context) or f"{user_id}@unknown.com"
        language = determine_language(user_mail)

        # 使用 HelpCardBuilder 輸出 Adaptive Card（與 app_bak 一致透過 Attachment）
        welcome_text = welcome_msg or {
            "zh-TW": "🛠️ 功能選單",
            "en": "🛠️ Function Menu",
            "ja": "🛠️ 機能メニュー",
        }.get(language, "🛠️ 功能選單")

        card = self.help_card_builder.build_help_card(language, welcome_text)
        await turn_context.send_activity(card)
